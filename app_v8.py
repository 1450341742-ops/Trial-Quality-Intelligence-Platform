"""Trial Quality Intelligence Platform V8

V8 focuses on engineering readiness:
- db_adapter healthcheck
- placeholder-based template rendering for docx/pptx (basic text replacement)
- delivery package ZIP generation
- notification settings placeholders for email/DingTalk
- Docker/default entry readiness
- quality gate dashboard

It reuses app_v4/app_v5/app_v6/app_v7 to keep compatibility.
"""

from __future__ import annotations

import json
import os
import zipfile
from datetime import datetime
from io import BytesIO
from pathlib import Path

import pandas as pd
import streamlit as st
from docx import Document
from pptx import Presentation

import app_v4 as core
import app_v5 as flow
import app_v6 as v6
import app_v7 as v7
import db_adapter

APP_TITLE = "Trial Quality Intelligence Platform V8"
RENDERED_DIR = Path(os.getenv("TQIP_RENDERED_DIR", "rendered_outputs"))
RENDERED_DIR.mkdir(exist_ok=True)

PLACEHOLDER_HELP = {
    "{{project_name}}": "项目名称",
    "{{sponsor_name}}": "申办方名称",
    "{{protocol_no}}": "方案编号",
    "{{indication}}": "适应症",
    "{{phase}}": "研究阶段",
    "{{readiness_score}}": "核查准备评分",
    "{{generated_at}}": "生成时间",
    "{{high_risk_count}}": "高风险问题数量",
    "{{open_capa_count}}": "未关闭CAPA数量",
}


def init_v8_db():
    v7.init_v7_db()
    conn = core.get_conn()
    cur = conn.cursor()
    cur.execute("""
        CREATE TABLE IF NOT EXISTS notification_settings (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            channel TEXT,
            target TEXT,
            enabled INTEGER DEFAULT 0,
            description TEXT,
            updated_by TEXT,
            updated_at TEXT
        )
    """)
    cur.execute("""
        CREATE TABLE IF NOT EXISTS rendered_documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            project_id INTEGER,
            template_file_id INTEGER,
            output_name TEXT,
            output_path TEXT,
            output_type TEXT,
            created_by TEXT,
            created_at TEXT
        )
    """)
    conn.commit()
    conn.close()


def login_screen():
    st.set_page_config(page_title=APP_TITLE, page_icon="🧬", layout="wide")
    core.css()
    core.hero("Trial Quality Intelligence Platform V8", "工程化增强版｜DB适配层｜模板占位符套版｜交付包ZIP｜提醒配置｜质量门禁")
    c1, c2, c3 = st.columns([1, 1.1, 1])
    with c2:
        st.markdown("<div class='card'>", unsafe_allow_html=True)
        st.subheader("登录")
        username = st.text_input("账号", value="admin")
        password = st.text_input("密码", type="password", value="admin123")
        if st.button("进入V8工程化增强版", use_container_width=True):
            user = core.authenticate(username, password)
            if user:
                st.session_state.update({"logged_in": True, "username": user["username"], "role": user["role"], "display_name": user["display_name"]})
                core.log("登录", "系统", "V8登录成功")
                st.rerun()
            else:
                st.error("账号或密码错误")
        st.info("默认账号：admin/admin123；qa/qa123；pm/pm123")
        st.markdown("</div>", unsafe_allow_html=True)


def project_context(project: dict) -> dict:
    findings = core.query_df("SELECT * FROM findings WHERE project_id=?", (project["id"],))
    score, _ = core.inspection_score(findings)
    high = int(findings["risk_level"].isin(["高风险", "极高风险"]).sum()) if not findings.empty else 0
    open_capa = int(findings["capa_status"].astype(str).str.contains("未|逾期|进行中|待", na=False).sum()) if not findings.empty else 0
    return {
        "project_name": project.get("project_name", ""),
        "sponsor_name": project.get("sponsor_name", ""),
        "protocol_no": project.get("protocol_no", ""),
        "indication": project.get("indication", ""),
        "phase": project.get("phase", ""),
        "readiness_score": str(score),
        "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "high_risk_count": str(high),
        "open_capa_count": str(open_capa),
    }


def replace_text(text: str, context: dict) -> str:
    out = text
    for key, value in context.items():
        out = out.replace("{{" + key + "}}", str(value))
    return out


def render_docx_template(template_path: str, context: dict) -> BytesIO:
    doc = Document(template_path)
    for p in doc.paragraphs:
        for run in p.runs:
            run.text = replace_text(run.text, context)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    for run in p.runs:
                        run.text = replace_text(run.text, context)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def render_pptx_template(template_path: str, context: dict) -> BytesIO:
    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = replace_text(run.text, context)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell.text = replace_text(cell.text, context)
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def render_v8_dashboard(projects: pd.DataFrame):
    core.hero("V8工程化驾驶舱", "数据库健康、质量门禁、模板套版、交付包和提醒配置总览")
    health = db_adapter.healthcheck()
    rendered = core.query_df("SELECT * FROM rendered_documents")
    packs = core.query_df("SELECT * FROM export_packs")
    tasks = core.query_df("SELECT * FROM tasks")
    overdue = v7.overdue_tasks(tasks)
    c1, c2, c3, c4, c5 = st.columns(5)
    with c1: core.metric_card("DB状态", health["status"], health["mode"])
    with c2: core.metric_card("项目数", len(projects), "Portfolio")
    with c3: core.metric_card("套版文档", len(rendered), "Rendered")
    with c4: core.metric_card("交付包", len(packs), "Packages")
    with c5: core.metric_card("逾期任务", len(overdue), "Overdue")
    st.subheader("数据库健康检查")
    st.json(health)
    st.subheader("质量门禁建议")
    gates = quality_gates()
    st.dataframe(gates, use_container_width=True)


def quality_gates() -> pd.DataFrame:
    checks = []
    health = db_adapter.healthcheck()
    checks.append({"门禁项": "数据库连接", "状态": "通过" if health["status"] == "ok" else "失败", "说明": health.get("error", "")})
    checks.append({"门禁项": "AI Key环境变量", "状态": "通过" if os.getenv("OPENAI_API_KEY") or os.getenv("DEEPSEEK_API_KEY") else "提醒", "说明": "未设置时仍可使用本地规则引擎"})
    checks.append({"门禁项": "模板上传目录", "状态": "通过" if Path("template_uploads").exists() else "提醒", "说明": "用于模板套版"})
    checks.append({"门禁项": "输出目录", "状态": "通过" if RENDERED_DIR.exists() else "失败", "说明": str(RENDERED_DIR)})
    checks.append({"门禁项": "默认账号", "状态": "提醒", "说明": "客户试点前请修改admin默认密码"})
    return pd.DataFrame(checks)


def render_template_rendering(project: dict):
    core.hero("模板占位符套版", "基于已上传模板文件进行基础占位符替换，支持docx/pptx")
    st.subheader("支持占位符")
    st.json(PLACEHOLDER_HELP)
    files = core.query_df("SELECT * FROM template_files WHERE status='启用' ORDER BY id DESC")
    if files.empty:
        st.info("请先在模板文件上传页面上传docx/pptx模板。")
        return
    labels = {f"{r.id}｜{r.template_name}｜{r.file_name}": int(r.id) for r in files.itertuples()}
    file_id = labels[st.selectbox("选择模板文件", list(labels.keys()))]
    row = files[files["id"] == file_id].iloc[0].to_dict()
    context = project_context(project)
    st.subheader("本次替换数据")
    st.json(context)
    if st.button("生成套版文件", use_container_width=True):
        path = row.get("file_path")
        suffix = Path(path).suffix.lower()
        try:
            if suffix == ".docx":
                buf = render_docx_template(path, context)
                mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                output_name = f"{project['project_name']}_{row.get('template_name')}_套版.docx"
            elif suffix == ".pptx":
                buf = render_pptx_template(path, context)
                mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                output_name = f"{project['project_name']}_{row.get('template_name')}_套版.pptx"
            else:
                st.error("当前仅支持docx/pptx模板套版。")
                return
            output_path = RENDERED_DIR / output_name
            output_path.write_bytes(buf.getvalue())
            core.execute(
                "INSERT INTO rendered_documents(project_id, template_file_id, output_name, output_path, output_type, created_by, created_at) VALUES (?, ?, ?, ?, ?, ?, ?)",
                (project["id"], file_id, output_name, str(output_path), suffix, st.session_state.get("username"), datetime.now().isoformat(timespec="seconds")),
            )
            core.log("模板套版", "rendered_documents", output_name)
            st.success("套版文件已生成。")
            st.download_button("下载套版文件", data=buf, file_name=output_name, mime=mime)
        except Exception as exc:
            st.error(f"套版失败：{exc}")
    st.subheader("已生成套版文档")
    rendered = core.query_df("SELECT * FROM rendered_documents WHERE project_id=? ORDER BY id DESC", (project["id"],))
    st.dataframe(rendered, use_container_width=True)


def render_zip_pack(project: dict):
    core.hero("交付包ZIP生成", "将已生成报告、套版文件和清单打包为ZIP，便于交付客户或内部归档")
    rendered = core.query_df("SELECT * FROM rendered_documents WHERE project_id=? ORDER BY id DESC", (project["id"],))
    packs = core.query_df("SELECT * FROM export_packs WHERE project_id=? ORDER BY id DESC", (project["id"],))
    st.subheader("已生成套版文件")
    st.dataframe(rendered, use_container_width=True)
    st.subheader("交付包元数据")
    st.dataframe(packs, use_container_width=True)
    include_rendered = st.checkbox("包含已生成套版文件", value=True)
    include_manifest = st.checkbox("包含交付清单manifest.json", value=True)
    if st.button("生成ZIP交付包", use_container_width=True):
        zip_buf = BytesIO()
        manifest = {"project": project, "generated_at": datetime.now().isoformat(timespec="seconds"), "files": []}
        with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
            if include_rendered and not rendered.empty:
                for _, r in rendered.iterrows():
                    path = Path(str(r.get("output_path")))
                    if path.exists():
                        arcname = f"rendered/{path.name}"
                        zf.write(path, arcname)
                        manifest["files"].append(arcname)
            if include_manifest:
                zf.writestr("manifest.json", json.dumps(manifest, ensure_ascii=False, indent=2, default=str))
            if not manifest["files"]:
                zf.writestr("README.txt", "当前交付包未包含套版文件。请先在模板占位符套版页面生成文件。")
        zip_buf.seek(0)
        output_name = f"{project['project_name']}_交付包_{datetime.now().strftime('%Y%m%d%H%M%S')}.zip"
        core.log("生成ZIP交付包", "export_pack_zip", output_name)
        st.download_button("下载ZIP交付包", zip_buf, file_name=output_name, mime="application/zip")


def render_notification_settings():
    core.hero("提醒配置", "钉钉/邮件提醒占位配置，用于后续任务逾期、复核SLA和核查倒计时通知")
    with st.form("notification_form_v8"):
        channel = st.selectbox("提醒渠道", ["DingTalk", "Email", "Webhook", "企业微信"])
        target = st.text_input("目标地址/机器人Webhook/邮箱")
        enabled = st.checkbox("启用", value=False)
        description = st.text_area("说明", value="用于任务逾期和复核SLA提醒", height=80)
        if st.form_submit_button("保存提醒配置"):
            core.execute(
                "INSERT INTO notification_settings(channel, target, enabled, description, updated_by, updated_at) VALUES (?, ?, ?, ?, ?, ?)",
                (channel, target, 1 if enabled else 0, description, st.session_state.get("username"), datetime.now().isoformat(timespec="seconds")),
            )
            core.log("保存提醒配置", "notification_settings", channel)
            st.success("提醒配置已保存。当前为占位配置，尚未实际发送。")
    st.dataframe(core.query_df("SELECT * FROM notification_settings ORDER BY id DESC"), use_container_width=True)
    st.info("V8仅保存提醒配置。V9可接入真实钉钉机器人、SMTP或企业微信Webhook。")


def render_engineering_center():
    core.hero("工程化中心", "测试、Docker、环境变量和运行入口检查")
    st.subheader("建议命令")
    st.code("pip install -r requirements.txt\npytest\nstreamlit run app_v8.py\ndocker compose up --build", language="bash")
    st.subheader("关键环境变量")
    st.code("DATABASE_URL=postgresql://...\nTQIP_DB_PATH=trial_quality_v8.db\nTQIP_UPLOAD_DIR=uploads\nTQIP_TEMPLATE_DIR=template_uploads\nTQIP_RENDERED_DIR=rendered_outputs\nOPENAI_API_KEY=...\nOPENAI_BASE_URL=...", language="bash")
    st.subheader("质量门禁")
    st.dataframe(quality_gates(), use_container_width=True)


def main():
    init_v8_db()
    if not st.session_state.get("logged_in"):
        login_screen()
        return
    st.set_page_config(page_title=APP_TITLE, page_icon="🧬", layout="wide")
    core.css()
    st.sidebar.markdown(f"### {st.session_state.get('display_name')}")
    st.sidebar.caption(f"角色：{st.session_state.get('role')}")
    if st.sidebar.button("退出登录"):
        core.log("退出登录", "系统", "用户退出")
        st.session_state.clear()
        st.rerun()

    menus = [
        "V8工程化驾驶舱", "数据库迁移中心", "客户中心", "客户-项目绑定", "客户级项目视图", "项目管理", "文件解析", "AI结构化解析", "AI字段映射", "问题清单", "风险分析", "专家复核", "复核转任务", "复核SLA中心", "ICF版本链核查", "SAE报告链核查", "中心文件评分", "中心文件批量导入", "任务逾期中心", "交付包生成", "模板中心", "模板文件上传", "模板占位符套版", "交付包ZIP", "提醒配置", "上线检查清单", "工程化中心", "用户管理", "修改密码", "报告中心", "数据治理", "系统设置"
    ]
    role = st.session_state.get("role")
    if role not in ["系统管理员", "申办方QA负责人"]:
        menus = [m for m in menus if m in ["V8工程化驾驶舱", "客户级项目视图", "项目管理", "文件解析", "问题清单", "风险分析", "专家复核", "复核SLA中心", "ICF版本链核查", "SAE报告链核查", "中心文件评分", "任务逾期中心", "模板占位符套版", "交付包ZIP", "报告中心", "修改密码"]]
    menu = st.sidebar.radio("功能导航", menus)
    project, projects = core.project_selector() if menu not in ["V8工程化驾驶舱", "数据库迁移中心", "客户中心", "客户-项目绑定", "客户级项目视图", "项目管理", "用户管理", "模板中心", "模板文件上传", "修改密码", "上线检查清单", "提醒配置", "工程化中心", "系统设置"] else (None, core.query_df("SELECT * FROM projects ORDER BY updated_at DESC,id DESC"))

    if menu == "V8工程化驾驶舱": render_v8_dashboard(projects)
    elif menu == "数据库迁移中心": v7.render_database_migration_center()
    elif menu == "客户中心": flow.render_client_center()
    elif menu == "客户-项目绑定": v6.render_project_company_binding()
    elif menu == "客户级项目视图": v7.render_customer_project_view()
    elif menu == "项目管理": core.render_project_management(projects)
    elif menu == "用户管理": flow.render_user_admin()
    elif menu == "模板中心": flow.render_template_center()
    elif menu == "模板文件上传": v7.render_template_file_upload()
    elif menu == "修改密码": v6.render_password_change()
    elif menu == "上线检查清单": v7.render_v7_release_checklist()
    elif menu == "提醒配置": render_notification_settings()
    elif menu == "工程化中心": render_engineering_center()
    elif menu == "系统设置": core.render_settings()
    else:
        if not project:
            st.info("请先创建并选择项目。")
            return
        if menu == "文件解析": core.render_file_parse(project)
        elif menu == "AI结构化解析": flow.render_ai_extraction_center(project)
        elif menu == "AI字段映射": v6.render_ai_field_mapping(project)
        elif menu == "问题清单": core.render_findings(project)
        elif menu == "风险分析": core.render_risk_analysis(project)
        elif menu == "专家复核": flow.render_expert_review(project)
        elif menu == "复核转任务": v6.render_review_to_task(project)
        elif menu == "复核SLA中心": v7.render_review_sla_center(project)
        elif menu == "ICF版本链核查": v6.render_icf_chain_check(project)
        elif menu == "SAE报告链核查": v6.render_sae_chain_check(project)
        elif menu == "中心文件评分": v6.render_center_file_score(project)
        elif menu == "中心文件批量导入": v7.render_center_file_batch_import(project)
        elif menu == "任务逾期中心": v7.render_task_overdue_center(project)
        elif menu == "交付包生成": v7.render_export_pack(project)
        elif menu == "模板占位符套版": render_template_rendering(project)
        elif menu == "交付包ZIP": render_zip_pack(project)
        elif menu == "报告中心": core.render_report(project)
        elif menu == "数据治理": core.render_data_governance(project)


if __name__ == "__main__":
    main()
