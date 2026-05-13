import hashlib
import json
import os
import sqlite3
from datetime import datetime, timedelta
from io import BytesIO
from pathlib import Path

import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt

try:
    import fitz
except Exception:
    fitz = None

try:
    from openai import OpenAI
except Exception:
    OpenAI = None

APP_TITLE = "Trial Quality Intelligence Platform V4"
DB_PATH = Path(os.getenv("TQIP_DB_PATH", "trial_quality_v4.db"))
UPLOAD_DIR = Path(os.getenv("TQIP_UPLOAD_DIR", "uploads"))
UPLOAD_DIR.mkdir(exist_ok=True)

DEFAULT_USERS = [
    ("admin", "admin123", "系统管理员", "系统管理员"),
    ("qa", "qa123", "申办方QA负责人", "QA负责人"),
    ("pm", "pm123", "项目经理PM", "项目经理"),
]

ROLE_PERMISSIONS = {
    "系统管理员": ["全部"],
    "申办方QA负责人": ["全部"],
    "项目经理PM": ["执行驾驶舱", "项目管理", "文件解析", "问题清单", "风险分析", "CAPA中心", "任务中心", "报告中心", "数据治理"],
    "注册负责人": ["管理层驾驶舱", "风险分析", "核查问答", "证据矩阵", "报告中心"],
    "只读用户": ["管理层驾驶舱", "执行驾驶舱", "报告中心"],
}

FINDING_CATEGORIES = {
    "知情同意": ["知情", "ICF", "同意书", "签署", "版本"],
    "入选/排除标准": ["入选", "排除", "纳入", "标准", "筛选失败"],
    "AE/SAE": ["AE", "SAE", "不良事件", "严重不良事件", "妊娠", "住院"],
    "方案偏离": ["方案偏离", "违背", "超窗", "访视窗", "未按方案"],
    "试验用药品管理": ["药品", "IP", "给药", "回收", "温度", "发放"],
    "数据完整性": ["EDC", "原始", "源数据", "一致", "缺失", "修改痕迹"],
    "主要终点": ["主要终点", "终点", "疗效评价", "影像", "阅片"],
    "研究者文件夹": ["研究者文件夹", "ISF", "授权表", "培训记录", "伦理批件", "授权分工"],
    "伦理/文件版本": ["伦理", "批件", "修正案", "版本", "递交"],
    "供应商管理": ["CRO", "SMO", "供应商", "中心实验室", "影像供应商"],
}

CRITICAL_DOMAINS = ["知情同意", "AE/SAE", "数据完整性", "主要终点", "伦理/文件版本", "入选/排除标准"]
LOW_QUALITY_CAPA_PHRASES = ["加强培训", "加强管理", "后续注意", "已整改", "已知晓", "CRA已提醒", "研究者已知晓"]


def css():
    st.markdown("""
    <style>
    .stApp { background: linear-gradient(180deg,#F8FBFF 0%,#FFFFFF 34%); }
    section[data-testid="stSidebar"] { background:#08111F; }
    section[data-testid="stSidebar"] * { color:#E5E7EB !important; }
    .hero { padding:26px 30px;border-radius:24px;background:linear-gradient(135deg,#0F172A 0%,#1D4ED8 55%,#06B6D4 100%);color:white;box-shadow:0 18px 42px rgba(37,99,235,.22);margin-bottom:20px; }
    .hero h1 { margin:0;font-size:34px; }
    .hero p { margin:10px 0 0 0;color:#DBEAFE;font-size:15px; }
    .card { background:white;padding:18px 20px;border-radius:18px;border:1px solid #E5E7EB;box-shadow:0 10px 28px rgba(15,23,42,.06); }
    .metric-card { background:white;padding:18px;border-radius:18px;border:1px solid #E5E7EB;box-shadow:0 10px 26px rgba(15,23,42,.06); }
    .metric-label { color:#64748B;font-size:13px;margin-bottom:8px; }
    .metric-value { color:#0F172A;font-weight:800;font-size:28px; }
    .metric-sub { color:#64748B;font-size:12px;margin-top:4px; }
    .section-title { font-size:22px;font-weight:800;color:#0F172A;margin:16px 0 10px; }
    </style>
    """, unsafe_allow_html=True)


def hero(title, subtitle):
    st.markdown(f"<div class='hero'><h1>{title}</h1><p>{subtitle}</p></div>", unsafe_allow_html=True)


def metric_card(label, value, sub=""):
    st.markdown(f"<div class='metric-card'><div class='metric-label'>{label}</div><div class='metric-value'>{value}</div><div class='metric-sub'>{sub}</div></div>", unsafe_allow_html=True)


def section(title):
    st.markdown(f"<div class='section-title'>{title}</div>", unsafe_allow_html=True)


def hash_password(password):
    return hashlib.sha256(password.encode("utf-8")).hexdigest()


def get_conn():
    return sqlite3.connect(DB_PATH, check_same_thread=False)


def execute(sql, params=()):
    conn = get_conn()
    cur = conn.cursor()
    cur.execute(sql, params)
    conn.commit()
    last_id = cur.lastrowid
    conn.close()
    return last_id


def query_df(sql, params=()):
    conn = get_conn()
    df = pd.read_sql_query(sql, conn, params=params)
    conn.close()
    return df


def init_db():
    conn = get_conn()
    cur = conn.cursor()
    cur.execute("""CREATE TABLE IF NOT EXISTS users (id INTEGER PRIMARY KEY AUTOINCREMENT, username TEXT UNIQUE, password_hash TEXT, role TEXT, display_name TEXT, status TEXT DEFAULT '启用', created_at TEXT)""")
    cur.execute("""CREATE TABLE IF NOT EXISTS projects (id INTEGER PRIMARY KEY AUTOINCREMENT, project_name TEXT NOT NULL, sponsor_name TEXT, protocol_no TEXT, indication TEXT, phase TEXT, planned_subjects INTEGER DEFAULT 0, actual_subjects INTEGER DEFAULT 0, site_count INTEGER DEFAULT 0, pm_name TEXT, qa_name TEXT, cro_name TEXT, smo_name TEXT, expected_inspection_date TEXT, project_status TEXT DEFAULT '进行中', created_at TEXT, updated_at TEXT)""")
    cur.execute("""CREATE TABLE IF NOT EXISTS files (id INTEGER PRIMARY KEY AUTOINCREMENT, project_id INTEGER, file_name TEXT, document_type TEXT, file_path TEXT, extracted_text TEXT, parse_summary TEXT, created_by TEXT, created_at TEXT)""")
    cur.execute("""CREATE TABLE IF NOT EXISTS findings (id INTEGER PRIMARY KEY AUTOINCREMENT, project_id INTEGER, site_no TEXT, site_name TEXT, subject_no TEXT, category TEXT, severity TEXT, description TEXT, basis TEXT, capa TEXT, capa_status TEXT, risk_score INTEGER, risk_level TEXT, ai_suggestion TEXT, evidence_gap TEXT, fingerprint TEXT UNIQUE, created_at TEXT)""")
    cur.execute("""CREATE TABLE IF NOT EXISTS tasks (id INTEGER PRIMARY KEY AUTOINCREMENT, project_id INTEGER, task_name TEXT, priority TEXT, owner TEXT, due_date TEXT, status TEXT DEFAULT '未开始', source TEXT, created_at TEXT)""")
    cur.execute("""CREATE TABLE IF NOT EXISTS audit_logs (id INTEGER PRIMARY KEY AUTOINCREMENT, username TEXT, action TEXT, target TEXT, detail TEXT, created_at TEXT)""")
    cur.execute("""CREATE TABLE IF NOT EXISTS ai_settings (id INTEGER PRIMARY KEY AUTOINCREMENT, provider TEXT, model_name TEXT, api_key_hint TEXT, enabled INTEGER DEFAULT 0, updated_at TEXT)""")
    conn.commit()
    conn.close()
    seed_users()


def seed_users():
    df = query_df("SELECT username FROM users")
    names = set(df["username"].tolist()) if not df.empty else set()
    for username, password, role, display_name in DEFAULT_USERS:
        if username not in names:
            execute("INSERT INTO users(username,password_hash,role,display_name,created_at) VALUES (?,?,?,?,?)", (username, hash_password(password), role, display_name, datetime.now().isoformat(timespec="seconds")))


def log(action, target="", detail=""):
    execute("INSERT INTO audit_logs(username,action,target,detail,created_at) VALUES (?,?,?,?,?)", (st.session_state.get("username", "anonymous"), action, target, detail, datetime.now().isoformat(timespec="seconds")))


def authenticate(username, password):
    df = query_df("SELECT * FROM users WHERE username=? AND status='启用'", (username,))
    if df.empty:
        return None
    row = df.iloc[0].to_dict()
    return row if row["password_hash"] == hash_password(password) else None


def has_permission(menu):
    role = st.session_state.get("role", "只读用户")
    perms = ROLE_PERMISSIONS.get(role, [])
    return "全部" in perms or menu in perms


def login_screen():
    st.set_page_config(page_title=APP_TITLE, page_icon="🧬", layout="wide")
    css()
    hero("Trial Quality Intelligence Platform V4", "商业化增强版｜AI解析｜中心评分｜受试者画像｜PPT汇报｜数据治理｜Docker部署")
    c1, c2, c3 = st.columns([1, 1.1, 1])
    with c2:
        st.markdown("<div class='card'>", unsafe_allow_html=True)
        st.subheader("登录")
        username = st.text_input("账号", value="admin")
        password = st.text_input("密码", type="password", value="admin123")
        if st.button("进入V4商业化驾驶舱", use_container_width=True):
            user = authenticate(username, password)
            if user:
                st.session_state.update({"logged_in": True, "username": user["username"], "role": user["role"], "display_name": user["display_name"]})
                log("登录", "系统", "V4登录成功")
                st.rerun()
            else:
                st.error("账号或密码错误")
        st.info("默认账号：admin/admin123；qa/qa123；pm/pm123")
        st.markdown("</div>", unsafe_allow_html=True)


def classify_category(text):
    value = str(text or "")
    for cat, keys in FINDING_CATEGORIES.items():
        if any(k.lower() in value.lower() for k in keys):
            return cat
    return "其他"


def normalize_severity(text):
    value = str(text or "")
    if any(k in value for k in ["严重", "Critical", "critical"]):
        return "严重问题"
    if any(k in value for k in ["主要", "Major", "major"]):
        return "主要问题"
    if any(k in value for k in ["建议", "Recommendation", "建议项"]):
        return "建议项"
    return "一般问题"


def risk_score_for_text(text, severity):
    score = {"严重问题": 10, "主要问题": 5, "一般问题": 2, "建议项": 1}.get(severity, 2)
    reasons = []
    rules = [
        (["受试者安全", "SAE", "严重不良事件", "死亡", "住院"], 8, "涉及受试者安全或SAE"),
        (["数据完整性", "EDC", "原始记录", "源数据", "一致性", "缺失"], 8, "涉及数据完整性或源数据溯源"),
        (["主要终点", "终点", "影像评价", "疗效评价"], 10, "涉及主要终点或关键疗效数据"),
        (["知情", "ICF", "同意书"], 6, "涉及知情同意合规"),
        (["入选", "排除", "入排"], 6, "涉及入排标准"),
        (["伦理", "批件", "修正案", "版本"], 5, "涉及伦理或版本链"),
        (["逾期", "未关闭", "未完成", "待补充"], 4, "存在未关闭或逾期风险"),
    ]
    for keys, weight, reason in rules:
        if any(k in str(text) for k in keys):
            score += weight
            reasons.append(reason)
    if score >= 25:
        return score, "极高风险", "；".join(reasons)
    if score >= 15:
        return score, "高风险", "；".join(reasons)
    if score >= 7:
        return score, "中风险", "；".join(reasons)
    return score, "低风险", "；".join(reasons) or "未触发高权重风险规则"


def evidence_gap_for_text(text, category):
    req = {
        "知情同意": ["伦理批件", "ICF版本清单", "签署页", "筛选检查时间线"],
        "伦理/文件版本": ["伦理批件", "版本递交记录", "启用日期", "中心执行版本清单"],
        "AE/SAE": ["原始病历", "AE/SAE表", "上报邮件/系统截图", "医学判断记录"],
        "数据完整性": ["原始记录", "EDC截图", "Query记录", "稽查追踪记录"],
        "主要终点": ["原始记录", "EDC截图", "终点评价表", "阅片/评价证据"],
        "试验用药品管理": ["药品台账", "温度记录", "发放回收记录", "偏差评估"],
        "入选/排除标准": ["筛选检查报告", "入排判断记录", "研究者确认", "随机化记录"],
    }.get(category, ["问题原始证据", "CAPA证据", "QA复核记录"])
    missing = [r for r in req if r not in str(text)]
    return "需准备/核实：" + "、".join(missing[:5]) if missing else "证据链较完整，建议QA复核确认。"


def capa_review(text):
    text = str(text or "").strip()
    issues = []
    if not text:
        issues.append("CAPA为空，无法判断是否可关闭。")
    for phrase in LOW_QUALITY_CAPA_PHRASES:
        if phrase in text:
            issues.append(f"存在低质量表述：{phrase}，需要补充具体对象、动作、证据和验证方式。")
    checks = [("根因", "未体现明确根因分析。"), ("证据", "未明确完成证据。"), ("预防", "未体现预防措施。"), ("有效性", "未体现CAPA有效性验证方式。")]
    for key, msg in checks:
        if key not in text and not (key == "证据" and any(k in text for k in ["记录", "截图", "文件", "复核"])):
            issues.append(msg)
    return {"score": max(0, 100 - len(issues) * 15), "decision": "建议补充后再关闭" if issues else "可进入QA复核关闭", "issues": issues or ["CAPA内容较完整，建议由QA结合证据文件最终确认。"]}


def call_ai(prompt, system="你是临床试验质量风险与核查准备专家。"):
    settings = query_df("SELECT * FROM ai_settings WHERE enabled=1 ORDER BY id DESC LIMIT 1")
    if settings.empty:
        return None, "未启用AI接口，当前使用本地规则引擎。"
    if OpenAI is None:
        return None, "未安装openai依赖，无法调用AI接口。"
    api_key = os.getenv("OPENAI_API_KEY") or os.getenv("DEEPSEEK_API_KEY")
    if not api_key:
        return None, "未设置OPENAI_API_KEY或DEEPSEEK_API_KEY环境变量。"
    row = settings.iloc[0].to_dict()
    base_url = os.getenv("OPENAI_BASE_URL")
    client = OpenAI(api_key=api_key, base_url=base_url) if base_url else OpenAI(api_key=api_key)
    try:
        resp = client.chat.completions.create(
            model=row.get("model_name") or "gpt-4.1-mini",
            messages=[{"role": "system", "content": system}, {"role": "user", "content": prompt}],
            temperature=0.2,
        )
        return resp.choices[0].message.content, "AI调用成功"
    except Exception as exc:
        return None, f"AI调用失败：{exc}"


def read_text_from_upload(uploaded_file):
    name = uploaded_file.name.lower()
    data = uploaded_file.getvalue()
    if name.endswith((".txt", ".md")):
        return data.decode("utf-8", errors="ignore")
    if name.endswith(".docx"):
        doc = Document(BytesIO(data))
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(parts)
    if name.endswith(".pdf") and fitz is not None:
        pdf = fitz.open(stream=data, filetype="pdf")
        return "\n".join([f"\n--- Page {i+1} ---\n{page.get_text('text')}" for i, page in enumerate(pdf)])
    return ""


def protocol_risk_parse(text):
    risk_rules = [
        ("入排标准", ["入选标准", "排除标准", "纳入标准"], "逐条核查受试者是否满足入排，重点关注筛选检查日期与入组日期。"),
        ("主要终点", ["主要终点", "primary endpoint", "终点"], "确认主要终点数据来源、评价时间点、原始记录与EDC一致性。"),
        ("AE/SAE", ["SAE", "严重不良事件", "不良事件", "AE"], "核查AE/SAE识别、记录、报告时限和医学判断链条。"),
        ("知情同意", ["知情同意", "ICF", "同意书"], "核查版本、签署日期、签署完整性及筛选前签署要求。"),
        ("伦理与版本链", ["伦理", "修正案", "版本", "批件", "递交"], "核查伦理递交、批件、方案/ICF版本和中心执行版本一致性。"),
        ("随机化/盲态", ["随机", "盲态", "设盲", "IWRS"], "核查随机分配、盲态保持和紧急揭盲记录。"),
        ("试验用药品", ["试验用药", "给药", "回收", "温度"], "核查药品接收、储存、发放、回收、温控和账物一致性。"),
    ]
    rows = []
    for item, keys, suggestion in risk_rules:
        matched = [k for k in keys if k.lower() in text.lower()]
        if matched:
            rows.append({"风险主题": item, "命中关键词": "、".join(matched), "风险等级": "高风险" if item in ["主要终点", "AE/SAE", "知情同意", "伦理与版本链"] else "中风险", "重点关注": suggestion})
    return pd.DataFrame(rows or [{"风险主题": "待人工确认", "命中关键词": "未命中核心关键词", "风险等级": "待确认", "重点关注": "请上传完整方案正文或方案摘要。"}])


def save_uploaded_file(uploaded_file, project_id, doc_type, text):
    fp = UPLOAD_DIR / f"{project_id}_{datetime.now().strftime('%Y%m%d%H%M%S')}_{uploaded_file.name.replace('/', '_')}"
    fp.write_bytes(uploaded_file.getbuffer())
    execute("INSERT INTO files(project_id,file_name,document_type,file_path,extracted_text,parse_summary,created_by,created_at) VALUES (?,?,?,?,?,?,?,?)", (project_id, uploaded_file.name, doc_type, str(fp), text[:30000], "已提取文本" if text else "已上传，未提取文本", st.session_state.get("username", ""), datetime.now().isoformat(timespec="seconds")))
    log("上传文件", "files", uploaded_file.name)


def fingerprint(project_id, site_no, subject_no, category, description):
    raw = f"{project_id}|{site_no}|{subject_no}|{category}|{str(description)[:120]}"
    return hashlib.md5(raw.encode("utf-8")).hexdigest()


def normalize_uploaded_findings(df, project_id):
    col_map = {"中心编号": "site_no", "中心名称": "site_name", "受试者编号": "subject_no", "问题分类": "category", "严重程度": "severity", "问题描述": "description", "依据": "basis", "CAPA": "capa", "整改状态": "capa_status"}
    norm = pd.DataFrame()
    for zh, en in col_map.items():
        norm[en] = df[zh] if zh in df.columns else ""
    inserted, skipped = 0, 0
    for idx, row in norm.iterrows():
        combined = " ".join(str(row.get(c, "")) for c in norm.columns)
        if not row["category"]:
            norm.at[idx, "category"] = classify_category(combined)
        norm.at[idx, "severity"] = normalize_severity(row["severity"])
        score, level, reason = risk_score_for_text(combined, norm.at[idx, "severity"])
        gap = evidence_gap_for_text(combined, norm.at[idx, "category"])
        fp = fingerprint(project_id, norm.at[idx, "site_no"], norm.at[idx, "subject_no"], norm.at[idx, "category"], norm.at[idx, "description"])
        norm.at[idx, "risk_score"] = score
        norm.at[idx, "risk_level"] = level
        norm.at[idx, "ai_suggestion"] = reason
        norm.at[idx, "evidence_gap"] = gap
        norm.at[idx, "fingerprint"] = fp
        try:
            execute("""INSERT INTO findings(project_id,site_no,site_name,subject_no,category,severity,description,basis,capa,capa_status,risk_score,risk_level,ai_suggestion,evidence_gap,fingerprint,created_at) VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)""", (project_id, str(norm.at[idx, "site_no"]), str(norm.at[idx, "site_name"]), str(norm.at[idx, "subject_no"]), str(norm.at[idx, "category"]), str(norm.at[idx, "severity"]), str(norm.at[idx, "description"]), str(norm.at[idx, "basis"]), str(norm.at[idx, "capa"]), str(norm.at[idx, "capa_status"]), int(score), str(level), reason, gap, fp, datetime.now().isoformat(timespec="seconds")))
            inserted += 1
        except sqlite3.IntegrityError:
            skipped += 1
    log("导入问题清单", "findings", f"新增{inserted}条，跳过重复{skipped}条")
    return norm, inserted, skipped


def inspection_score(findings):
    if findings.empty:
        return 70, ["尚未上传稽查问题清单，核查准备评分仅为初步估算。"]
    score = 100
    severe = int((findings["severity"] == "严重问题").sum())
    high = int(findings["risk_level"].isin(["高风险", "极高风险"]).sum())
    open_capa = int(findings["capa_status"].astype(str).str.contains("未|逾期|进行中|待", na=False).sum())
    endpoint = int(findings["category"].astype(str).str.contains("主要终点|数据完整性|AE/SAE|知情|伦理", na=False).sum())
    evidence = int(findings["evidence_gap"].astype(str).str.contains("需准备|核实", na=False).sum()) if "evidence_gap" in findings.columns else 0
    score -= severe * 10 + high * 5 + open_capa * 4 + endpoint * 3 + min(evidence, 10) * 2
    gaps = []
    if severe: gaps.append(f"存在 {severe} 项严重问题，需形成专项解释和补救证据。")
    if high: gaps.append(f"存在 {high} 项高风险/极高风险问题，建议核查前完成QA复核。")
    if open_capa: gaps.append(f"存在 {open_capa} 项CAPA未完全关闭或状态不清。")
    if endpoint: gaps.append(f"存在 {endpoint} 项涉及主要终点、AE/SAE、ICF、伦理或数据完整性的重点问题。")
    if evidence: gaps.append(f"存在 {evidence} 项问题证据链需要补充或核实。")
    return max(score, 0), gaps or ["当前未发现明显核查准备缺口，仍需结合原始文件确认。"]


def site_readiness(findings):
    if findings.empty:
        return pd.DataFrame()
    rows = []
    for (site_no, site_name), df in findings.groupby(["site_no", "site_name"], dropna=False):
        score, gaps = inspection_score(df)
        rows.append({"中心编号": site_no, "中心名称": site_name, "中心核查准备评分": score, "问题数量": len(df), "高风险问题": int(df["risk_level"].isin(["高风险", "极高风险"]).sum()), "未关闭CAPA": int(df["capa_status"].astype(str).str.contains("未|逾期|进行中|待", na=False).sum()), "主要缺口": "；".join(gaps[:2])})
    return pd.DataFrame(rows).sort_values("中心核查准备评分")


def subject_profile(findings):
    if findings.empty or "subject_no" not in findings.columns:
        return pd.DataFrame()
    rows = []
    for sub, df in findings.groupby("subject_no", dropna=False):
        if not str(sub).strip():
            continue
        rows.append({"受试者编号": sub, "涉及中心": ",".join(df["site_name"].dropna().astype(str).unique()), "问题数量": len(df), "风险分": int(df["risk_score"].sum()), "高风险问题": int(df["risk_level"].isin(["高风险", "极高风险"]).sum()), "涉及领域": "、".join(df["category"].dropna().astype(str).unique())})
    return pd.DataFrame(rows).sort_values("风险分", ascending=False) if rows else pd.DataFrame()


def project_selector():
    projects = query_df("SELECT * FROM projects ORDER BY updated_at DESC, id DESC")
    if projects.empty:
        st.sidebar.warning("请先创建项目")
        return None, projects
    labels = {f"{r.project_name}｜{r.protocol_no or '无方案编号'}": int(r.id) for r in projects.itertuples()}
    sel = st.sidebar.selectbox("选择项目", list(labels.keys()))
    return projects[projects["id"] == labels[sel]].iloc[0].to_dict(), projects


def render_management(projects):
    hero("管理层驾驶舱 V4", "商业化版本：项目组合风险、中心准备评分、CAPA缺口、证据链和汇报导出")
    findings = query_df("SELECT * FROM findings")
    tasks = query_df("SELECT * FROM tasks")
    score, gaps = inspection_score(findings)
    c1, c2, c3, c4, c5 = st.columns(5)
    with c1: metric_card("项目总数", len(projects), "Portfolio")
    with c2: metric_card("总核查准备评分", score, "Inspection Readiness")
    with c3: metric_card("高风险问题", int(findings["risk_level"].isin(["高风险", "极高风险"]).sum()) if not findings.empty else 0, "需优先复核")
    with c4: metric_card("CAPA待关闭", int(findings["capa_status"].astype(str).str.contains("未|逾期|进行中|待", na=False).sum()) if not findings.empty else 0, "影响核查")
    with c5: metric_card("待办任务", len(tasks), "Action Items")
    section("管理层结论")
    if score >= 85: st.success("整体状态较好，可进入核查前精细复核阶段。")
    elif score >= 70: st.warning("存在质量缺口，建议优先处理高风险中心、严重问题和CAPA。")
    else: st.error("核查准备风险较高，建议启动专项质量补救和模拟核查。")
    for gap in gaps: st.write("- " + gap)
    if not findings.empty:
        left, right = st.columns([1.2, 1])
        with left:
            sr = site_readiness(findings)
            st.plotly_chart(px.bar(sr.head(12), x="中心名称", y="中心核查准备评分", hover_data=["问题数量", "高风险问题"], title="中心核查准备评分 TOP低分中心"), use_container_width=True)
        with right:
            st.plotly_chart(px.pie(findings.groupby("category").size().reset_index(name="数量"), names="category", values="数量", title="问题分类占比"), use_container_width=True)


def render_project_management(projects):
    hero("项目管理与数据维护", "创建、查看、维护项目主数据；V4支持项目状态编辑和基础数据治理")
    with st.form("new_project_v4"):
        c1, c2, c3 = st.columns(3)
        project_name = c1.text_input("项目名称", "示例项目：TQIP-004 商业化验证项目")
        sponsor_name = c2.text_input("申办方名称", "某创新药申办方")
        protocol_no = c3.text_input("方案编号", "TQIP-004")
        indication = c1.text_input("适应症", "肿瘤/免疫/慢病等")
        phase = c2.selectbox("研究阶段", ["I期", "II期", "III期", "IV期", "真实世界研究"])
        expected_date = c3.date_input("预计核查/申报前检查日期")
        planned = c1.number_input("计划入组例数", min_value=0, value=180)
        actual = c2.number_input("实际入组例数", min_value=0, value=0)
        sites = c3.number_input("中心数量", min_value=0, value=18)
        pm = c1.text_input("PM", "")
        qa = c2.text_input("QA负责人", "")
        cro = c3.text_input("CRO", "")
        if st.form_submit_button("保存项目", use_container_width=True):
            now = datetime.now().isoformat(timespec="seconds")
            execute("INSERT INTO projects(project_name,sponsor_name,protocol_no,indication,phase,planned_subjects,actual_subjects,site_count,pm_name,qa_name,cro_name,expected_inspection_date,created_at,updated_at) VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?)", (project_name, sponsor_name, protocol_no, indication, phase, planned, actual, sites, pm, qa, cro, expected_date.isoformat(), now, now))
            log("创建项目", "projects", project_name)
            st.success("项目已创建")
    section("项目列表")
    st.dataframe(projects, use_container_width=True)
    if not projects.empty:
        section("项目状态维护")
        pid = st.selectbox("选择要更新的项目ID", projects["id"].tolist())
        status = st.selectbox("项目状态", ["进行中", "核查准备中", "已完成", "暂停", "归档"])
        if st.button("更新项目状态"):
            execute("UPDATE projects SET project_status=?, updated_at=? WHERE id=?", (status, datetime.now().isoformat(timespec="seconds"), int(pid)))
            log("更新项目状态", "projects", f"{pid}->{status}")
            st.success("已更新项目状态")


def render_file_parse(project):
    hero("文件解析与AI深度提取", "支持本地规则解析，也可在系统设置中启用OpenAI/DeepSeek接口进行深度结构化提取")
    doc_type = st.selectbox("资料类型", ["临床试验方案", "方案修正案", "知情同意书", "伦理批件", "稽查报告", "监查报告", "供应商报告", "其他"])
    uploaded = st.file_uploader("上传 Word/PDF/TXT/MD 文件", type=["docx", "txt", "md", "pdf"])
    if uploaded:
        text = read_text_from_upload(uploaded)
        save_uploaded_file(uploaded, project["id"], doc_type, text)
        if not text:
            st.warning("未提取到文本。扫描PDF请后续接OCR。")
        else:
            risks = protocol_risk_parse(text)
            st.session_state["last_protocol_risks"] = risks
            st.success(f"已提取文本 {len(text)} 字符。")
            st.dataframe(risks, use_container_width=True)
            if st.button("尝试AI深度结构化解析"):
                ai, msg = call_ai("请基于以下临床试验方案文本，输出JSON：项目概况、入选标准、排除标准、主要终点、关键时间窗、AE/SAE要求、知情同意风险、稽查重点。文本：\n" + text[:12000])
                st.info(msg)
                if ai: st.text_area("AI解析结果", value=ai, height=360)
    files = query_df("SELECT id,file_name,document_type,parse_summary,created_by,created_at FROM files WHERE project_id=? ORDER BY id DESC", (project["id"],))
    st.dataframe(files, use_container_width=True)


def render_findings(project):
    hero("问题清单导入与去重", "自动分类、评分、识别证据缺口；V4增加重复导入保护")
    uploaded = st.file_uploader("上传 CSV 或 XLSX", type=["csv", "xlsx"])
    if uploaded:
        df = pd.read_csv(uploaded) if uploaded.name.lower().endswith(".csv") else pd.read_excel(uploaded)
        norm, inserted, skipped = normalize_uploaded_findings(df, project["id"])
        st.success(f"解析完成：新增 {inserted} 条，跳过重复 {skipped} 条。")
        st.dataframe(norm, use_container_width=True)
    findings = query_df("SELECT * FROM findings WHERE project_id=? ORDER BY risk_score DESC,id DESC", (project["id"],))
    st.dataframe(findings, use_container_width=True)


def render_risk_analysis(project):
    hero("高级风险分析", "中心级核查准备评分、受试者风险画像、风险热力图、高风险问题和证据缺口")
    findings = query_df("SELECT * FROM findings WHERE project_id=?", (project["id"],))
    if findings.empty:
        st.info("请先上传问题清单。")
        return
    tab1, tab2, tab3, tab4, tab5 = st.tabs(["中心准备评分", "受试者画像", "风险热力图", "高风险问题", "证据缺口"])
    with tab1:
        sr = site_readiness(findings)
        st.dataframe(sr, use_container_width=True)
        st.plotly_chart(px.bar(sr, x="中心名称", y="中心核查准备评分", hover_data=["问题数量", "高风险问题", "主要缺口"], title="中心级 Inspection Readiness Score"), use_container_width=True)
    with tab2:
        sp = subject_profile(findings)
        st.dataframe(sp, use_container_width=True)
        if not sp.empty:
            st.plotly_chart(px.bar(sp.head(20), x="受试者编号", y="风险分", hover_data=["涉及中心", "涉及领域"], title="受试者风险画像 TOP20"), use_container_width=True)
    with tab3:
        heat = findings.pivot_table(index="site_name", columns="category", values="risk_score", aggfunc="sum", fill_value=0)
        st.plotly_chart(px.imshow(heat, text_auto=True, aspect="auto", title="中心 × 问题分类 风险热力图"), use_container_width=True)
    with tab4:
        st.dataframe(findings[findings["risk_level"].isin(["高风险", "极高风险"])].sort_values("risk_score", ascending=False), use_container_width=True)
    with tab5:
        st.dataframe(findings[["site_name", "subject_no", "category", "description", "evidence_gap", "risk_level"]], use_container_width=True)


def generate_qa(project, findings):
    rows = [{"角色": "申办方QA", "问题": f"请说明{project.get('project_name','本项目')}如何进行质量风险管理？", "建议回答": "申办方基于方案关键风险、中心发现、CAPA关闭情况和核查准备评分进行分层管理。", "需准备证据": "质量管理计划、稽查计划、风险评估表、CAPA跟踪表、项目质量分析报告"}]
    if not findings.empty:
        for _, row in findings[findings["risk_level"].isin(["高风险", "极高风险"])].head(10).iterrows():
            rows.append({"角色": "研究者/中心人员", "问题": f"中心发现{row.get('category','')}问题：{str(row.get('description',''))[:70]}，请解释原因和整改情况。", "建议回答": "说明发生原因、影响评估、纠正措施、预防措施和证据链。", "需准备证据": str(row.get("evidence_gap", "原始记录、EDC截图、CAPA证据"))})
    return pd.DataFrame(rows)


def render_capa(project):
    hero("CAPA中心", "审核CAPA质量，识别根因、证据、预防措施和有效性验证缺口")
    findings = query_df("SELECT * FROM findings WHERE project_id=? ORDER BY risk_score DESC", (project["id"],))
    if findings.empty:
        st.info("请先上传问题清单。")
        return
    selected = st.selectbox("选择问题", findings["id"].astype(str) + "｜" + findings["description"].astype(str).str.slice(0, 70))
    row = findings[findings["id"] == int(selected.split("｜")[0])].iloc[0]
    st.write("问题描述：", row["description"])
    result = capa_review(st.text_area("CAPA内容", value=str(row.get("capa", "")), height=180))
    c1, c2 = st.columns(2)
    with c1: metric_card("CAPA质量评分", result["score"], result["decision"])
    with c2: metric_card("风险等级", row["risk_level"], row.get("evidence_gap", ""))
    for issue in result["issues"]: st.warning(issue)


def generate_tasks(project_id, findings):
    execute("DELETE FROM tasks WHERE project_id=? AND source='系统自动生成'", (project_id,))
    _, gaps = inspection_score(findings)
    tasks = []
    due7 = (datetime.now() + timedelta(days=7)).date().isoformat()
    due15 = (datetime.now() + timedelta(days=15)).date().isoformat()
    for gap in gaps: tasks.append(("核查准备缺口补救：" + gap[:42], "高", "QA负责人", due7))
    if not findings.empty:
        for _, row in findings[findings["risk_level"].isin(["高风险", "极高风险"])].head(12).iterrows(): tasks.append((f"复核高风险问题：{row.get('site_name','')}｜{row.get('category','')}", "高", "QA负责人", due7))
        for _, row in findings[findings["capa_status"].astype(str).str.contains("未|逾期|进行中|待", na=False)].head(12).iterrows(): tasks.append((f"跟进CAPA关闭：{row.get('site_name','')}｜{str(row.get('description',''))[:22]}", "中", "PM/CRA", due15))
    for t in tasks:
        execute("INSERT INTO tasks(project_id,task_name,priority,owner,due_date,status,source,created_at) VALUES (?,?,?,?,?,?,?,?)", (project_id, t[0], t[1], t[2], t[3], "未开始", "系统自动生成", datetime.now().isoformat(timespec="seconds")))
    log("生成任务清单", "tasks", f"生成{len(tasks)}项任务")
    return len(tasks)


def render_tasks(project):
    hero("任务中心", "将风险缺口转化为核查前可执行行动项")
    findings = query_df("SELECT * FROM findings WHERE project_id=?", (project["id"],))
    if st.button("基于当前风险生成任务清单", use_container_width=True): st.success(f"已生成 {generate_tasks(project['id'], findings)} 项任务。")
    tasks = query_df("SELECT * FROM tasks WHERE project_id=? ORDER BY due_date ASC,id DESC", (project["id"],))
    st.dataframe(tasks, use_container_width=True)
    if not tasks.empty: st.download_button("下载任务清单CSV", tasks.to_csv(index=False).encode("utf-8-sig"), file_name=f"{project['project_name']}_核查前任务清单.csv", mime="text/csv")


def add_table(doc, df):
    if df.empty:
        doc.add_paragraph("暂无数据。")
        return
    table = doc.add_table(rows=1, cols=len(df.columns))
    for i, col in enumerate(df.columns): table.rows[0].cells[i].text = str(col)
    for _, row in df.iterrows():
        cells = table.add_row().cells
        for i, col in enumerate(df.columns): cells[i].text = str(row.get(col, ""))[:800]


def generate_word(project, findings, qa):
    doc = Document(); doc.styles["Normal"].font.name = "Arial"; doc.styles["Normal"].font.size = Pt(10)
    score, gaps = inspection_score(findings)
    doc.add_heading("项目质量风险与核查准备评估报告", level=1)
    for label, key in [("项目名称", "project_name"), ("申办方", "sponsor_name"), ("方案编号", "protocol_no"), ("适应症", "indication")]: doc.add_paragraph(f"{label}：{project.get(key,'')}")
    doc.add_paragraph(f"核查准备评分：{score}分")
    doc.add_heading("一、管理层结论", level=2)
    for gap in gaps: doc.add_paragraph(gap)
    if not findings.empty:
        doc.add_heading("二、中心核查准备评分", level=2); add_table(doc, site_readiness(findings))
        doc.add_heading("三、受试者风险画像", level=2); add_table(doc, subject_profile(findings).head(20))
        doc.add_heading("四、高风险问题清单", level=2); add_table(doc, findings[findings["risk_level"].isin(["高风险", "极高风险"])][["site_name", "subject_no", "category", "severity", "description", "risk_level", "evidence_gap"]].head(30))
    doc.add_heading("五、核查访谈问答", level=2); add_table(doc, qa)
    buf = BytesIO(); doc.save(buf); buf.seek(0); log("导出Word报告", "report", project.get("project_name", "")); return buf


def generate_ppt(project, findings):
    prs = Presentation()
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = "项目质量风险与核查准备汇报"
    title.placeholders[1].text = f"{project.get('project_name','')}｜{datetime.now().strftime('%Y-%m-%d')}"
    score, gaps = inspection_score(findings)
    slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = "管理层结论"
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.8), Inches(4.8)); tf = box.text_frame
    tf.text = f"核查准备评分：{score}分"
    for gap in gaps[:5]: p = tf.add_paragraph(); p.text = "• " + gap; p.font.size = PptPt(18)
    if not findings.empty:
        slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = "中心核查准备评分"
        sr = site_readiness(findings).head(8)
        table = slide.shapes.add_table(len(sr)+1, 4, Inches(0.4), Inches(1.2), Inches(9.2), Inches(4.8)).table
        headers = ["中心", "评分", "问题数", "高风险"]
        for i,h in enumerate(headers): table.cell(0,i).text = h
        for r, (_, row) in enumerate(sr.iterrows(), start=1):
            table.cell(r,0).text = str(row["中心名称"]); table.cell(r,1).text = str(row["中心核查准备评分"]); table.cell(r,2).text = str(row["问题数量"]); table.cell(r,3).text = str(row["高风险问题"])
        slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = "高风险问题摘要"
        high = findings[findings["risk_level"].isin(["高风险", "极高风险"])].head(8)
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(9), Inches(5)); tf = box.text_frame
        tf.text = "需优先复核的问题："
        for _, row in high.iterrows():
            p = tf.add_paragraph(); p.text = f"• {row.get('site_name','')}｜{row.get('category','')}｜{str(row.get('description',''))[:50]}"; p.font.size = PptPt(14)
    buf = BytesIO(); prs.save(buf); buf.seek(0); log("导出PPT报告", "report", project.get("project_name", "")); return buf


def render_report(project):
    hero("报告中心", "导出Word综合报告与PPT管理层汇报")
    findings = query_df("SELECT * FROM findings WHERE project_id=?", (project["id"],))
    qa = generate_qa(project, findings)
    score, gaps = inspection_score(findings)
    c1, c2 = st.columns(2)
    with c1: metric_card("核查准备评分", score, "写入Word/PPT")
    with c2: metric_card("导出类型", "Word + PPT", "商业演示版")
    for gap in gaps: st.write("- " + gap)
    c1, c2 = st.columns(2)
    with c1:
        if st.button("生成Word综合报告", use_container_width=True):
            st.download_button("下载Word报告", generate_word(project, findings, qa), file_name=f"{project['project_name']}_质量风险与核查准备报告.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    with c2:
        if st.button("生成PPT管理层汇报", use_container_width=True):
            st.download_button("下载PPT汇报", generate_ppt(project, findings), file_name=f"{project['project_name']}_质量风险管理层汇报.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")


def render_data_governance(project):
    hero("数据治理", "重复导入保护、项目数据清理、问题删除和演示数据维护")
    findings = query_df("SELECT id,site_name,subject_no,category,severity,description,risk_level,created_at FROM findings WHERE project_id=? ORDER BY id DESC", (project["id"],))
    st.dataframe(findings, use_container_width=True)
    if not findings.empty:
        del_ids = st.multiselect("选择要删除的问题ID", findings["id"].tolist())
        if st.button("删除所选问题") and del_ids:
            for fid in del_ids: execute("DELETE FROM findings WHERE id=?", (int(fid),))
            log("删除问题", "findings", f"删除{len(del_ids)}条")
            st.success("已删除所选问题")
    if st.button("清空当前项目任务"):
        execute("DELETE FROM tasks WHERE project_id=?", (project["id"],)); log("清空任务", "tasks", project.get("project_name", "")); st.success("已清空任务")


def render_settings():
    hero("系统设置", "AI接口、用户、日志和部署参数")
    tab1, tab2, tab3, tab4 = st.tabs(["AI接口", "用户", "操作日志", "部署信息"])
    with tab1:
        with st.form("ai_settings_v4"):
            provider = st.selectbox("模型供应商", ["OpenAI", "Azure OpenAI", "DeepSeek", "通义千问", "私有化模型"])
            model = st.text_input("模型名称", "gpt-4.1-mini")
            hint = st.text_input("API Key提示", "使用环境变量 OPENAI_API_KEY / DEEPSEEK_API_KEY")
            enabled = st.checkbox("启用AI接口", value=False)
            if st.form_submit_button("保存配置"):
                execute("INSERT INTO ai_settings(provider,model_name,api_key_hint,enabled,updated_at) VALUES (?,?,?,?,?)", (provider, model, hint, 1 if enabled else 0, datetime.now().isoformat(timespec="seconds")))
                log("保存AI配置", "ai_settings", provider); st.success("已保存AI配置")
        st.info("DeepSeek可通过设置OPENAI_BASE_URL为兼容地址进行调用。API Key不建议保存到数据库，请放在部署环境变量。")
    with tab2: st.dataframe(query_df("SELECT id,username,role,display_name,status,created_at FROM users"), use_container_width=True)
    with tab3: st.dataframe(query_df("SELECT * FROM audit_logs ORDER BY id DESC LIMIT 300"), use_container_width=True)
    with tab4:
        st.code(f"DB_PATH={DB_PATH}\nUPLOAD_DIR={UPLOAD_DIR}\nOPENAI_API_KEY={'已设置' if os.getenv('OPENAI_API_KEY') else '未设置'}\nOPENAI_BASE_URL={os.getenv('OPENAI_BASE_URL','未设置')}")


def main():
    init_db()
    if not st.session_state.get("logged_in"):
        login_screen(); return
    st.set_page_config(page_title=APP_TITLE, page_icon="🧬", layout="wide")
    css()
    st.sidebar.markdown(f"### {st.session_state.get('display_name')}")
    st.sidebar.caption(f"角色：{st.session_state.get('role')}")
    if st.sidebar.button("退出登录"):
        log("退出登录", "系统", "用户退出"); st.session_state.clear(); st.rerun()
    menus = ["管理层驾驶舱", "项目管理", "文件解析", "问题清单", "风险分析", "CAPA中心", "核查问答", "任务中心", "报告中心", "数据治理", "系统设置"]
    visible = [m for m in menus if has_permission(m)]
    menu = st.sidebar.radio("功能导航", visible)
    project, projects = project_selector() if menu not in ["管理层驾驶舱", "项目管理", "系统设置"] else (None, query_df("SELECT * FROM projects ORDER BY updated_at DESC,id DESC"))
    if menu == "管理层驾驶舱": render_management(projects)
    elif menu == "项目管理": render_project_management(projects)
    elif menu == "系统设置": render_settings()
    else:
        if not project: st.info("请先创建并选择项目。"); return
        if menu == "文件解析": render_file_parse(project)
        elif menu == "问题清单": render_findings(project)
        elif menu == "风险分析": render_risk_analysis(project)
        elif menu == "CAPA中心": render_capa(project)
        elif menu == "核查问答": hero("核查问答", "基于高风险问题生成访谈问题和证据准备清单"); st.dataframe(generate_qa(project, query_df("SELECT * FROM findings WHERE project_id=?", (project["id"],))), use_container_width=True)
        elif menu == "任务中心": render_tasks(project)
        elif menu == "报告中心": render_report(project)
        elif menu == "数据治理": render_data_governance(project)


if __name__ == "__main__":
    main()
